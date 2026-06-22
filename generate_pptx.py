"""
generate_pptx.py

Generates RTP_Analyzer_PoC_Achievements_Demo.pptx using python-pptx and matplotlib.

Usage:
  1. Create a Python virtual environment: python -m venv .venv
  2. Activate it: source .venv/bin/activate  (Linux/macOS) or .venv\Scripts\activate (Windows)
  3. Install dependencies: pip install python-pptx matplotlib
  4. Run: python generate_pptx.py

This script creates charts and a 12-slide inspirational deck with speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import os

# Metrics and content (from analysis)
METRICS = {
    'removed_triggers': 20,
    'complexity': {'Complex':1, 'High':6, 'Medium':40, 'Low':20},
    'endpoints': {'LOV':13, 'CRUD':6, 'ACTION':5, 'VALIDATION':4},
    'program_units_procedures': 3,
    'reports_called': 1,
    'ui_buttons_removed': 9,
    'fmb_before_kb': 700,
    'fmb_after_kb': 265,
    'security_issues_count': 3
}

OUT_PPTX = 'RTP_Analyzer_PoC_Achievements_Demo.pptx'
ASSETS_DIR = 'pptx_assets'
os.makedirs(ASSETS_DIR, exist_ok=True)

# Helper: create simple charts and save as PNG

def make_bar_trigger_complexity(path):
    labels = list(METRICS['complexity'].keys())
    values = [METRICS['complexity'][k] for k in labels]
    fig, ax = plt.subplots(figsize=(6,3))
    bars = ax.bar(labels, values, color=['#d9534f','#f0ad4e','#ffd24a','#5bc0de'])
    ax.set_title('Trigger Complexity Distribution')
    ax.set_ylabel('Count')
    for bar in bars:
        yval = bar.get_height()
        ax.text(bar.get_x()+bar.get_width()/2, yval + 0.3, int(yval), ha='center', va='bottom')
    fig.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def make_pie_endpoints(path):
    labels = list(METRICS['endpoints'].keys())
    sizes = [METRICS['endpoints'][k] for k in labels]
    colors = ['#6aa84f','#3c78d8','#f6b26b','#8e7cc3']
    fig, ax = plt.subplots(figsize=(4,4))
    ax.pie(sizes, labels=labels, autopct='%1.0f%%', colors=colors, startangle=140)
    ax.set_title('Endpoint Conversion Types')
    fig.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def make_bar_fmb_size(path):
    labels = ['Before (KB)', 'After (KB)']
    values = [METRICS['fmb_before_kb'], METRICS['fmb_after_kb']]
    fig, ax = plt.subplots(figsize=(4,3))
    bars = ax.bar(labels, values, color=['#444444','#2a9d8f'])
    ax.set_title('FMB Size Reduction')
    for bar in bars:
        yval = bar.get_height()
        ax.text(bar.get_x()+bar.get_width()/2, yval + 5, int(yval), ha='center', va='bottom')
    fig.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)

# Create assets
bar_complexity = os.path.join(ASSETS_DIR, 'bar_complexity.png')
pie_endpoints = os.path.join(ASSETS_DIR, 'pie_endpoints.png')
bar_fmb = os.path.join(ASSETS_DIR, 'bar_fmb.png')

make_bar_trigger_complexity(bar_complexity)
make_pie_endpoints(pie_endpoints)
make_bar_fmb_size(bar_fmb)

# Create presentation
prs = Presentation()
blank_layout = prs.slide_layouts[6]
title_layout = prs.slide_layouts[0]

# Utility to add text box

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=(0,0,0)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return tf

# 1 Title slide
slide = prs.slides.add_slide(title_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'RTP Analyzer PoC — Achievements & KPIs'
subtitle.text = 'Migration readiness, security findings, and next steps — Inspiring demo'
slide.notes_slide.notes_text_frame.text = 'Introduce the RTP Analyzer PoC, scope, and headline outcomes: endpoint mapping, size reduction, and security findings.'

# 2 Executive summary
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(1.5), 'Executive Summary', font_size=28, bold=True)
summary_text = (
    f"PoC analysis produced {sum(METRICS['endpoints'].values())} endpoint mappings, removed {METRICS['removed_triggers']} unused triggers, "
    f"reduced FMB size from {METRICS['fmb_before_kb']}KB to {METRICS['fmb_after_kb']}KB, and surfaced {METRICS['security_issues_count']} critical security issues."
)
add_textbox(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(2), summary_text, font_size=18)
slide.notes_slide.notes_text_frame.text = 'One-sentence overview and top outcomes to hook the audience.'

# 3 Scope & approach
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.8), 'Scope & Approach', font_size=28, bold=True)
scope_points = (
    'Analyzed Oracle Forms (FMB), triggers, program units, DB calls, and UI elements using the RTP Analyzer tool.\n'
    'Goals: remove dead code, classify complexity, map endpoints for BE, propose FE wireframes, and find security issues.'
)
add_textbox(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(2), scope_points, font_size=16)
slide.notes_slide.notes_text_frame.text = 'Describe the tool and the high-level method — static analysis and migration mapping.'

# 4 KPIs
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7), 'Key KPIs', font_size=28, bold=True)
kp_text = (
    f"- Code reduction: FMB size reduced by {int((1 - METRICS['fmb_after_kb']/METRICS['fmb_before_kb'])*100)}% ({METRICS['fmb_before_kb']}KB → {METRICS['fmb_after_kb']}KB)\n"
    f"- Triggers removed: {METRICS['removed_triggers']}\n"
    f"- Endpoint conversions identified: {sum(METRICS['endpoints'].values())} (LOV/CRUD/ACTION/VALIDATION)\n"
    f"- UI simplification: {METRICS['ui_buttons_removed']} buttons removed\n"
    f"- Business logic centralized: {METRICS['program_units_procedures']} program-unit procedures\n"
    f"- Security risks found: {METRICS['security_issues_count']} (priority remediation)"
)
add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(3.5), kp_text, font_size=16)
slide.notes_slide.notes_text_frame.text = 'Highlight measurable outcomes for stakeholders.'

# 5 Trigger complexity chart
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6), 'Trigger Complexity', font_size=28, bold=True)
slide.shapes.add_picture(bar_complexity, Inches(1), Inches(1), width=Inches(8))
slide.notes_slide.notes_text_frame.text = 'Discuss prioritization: Complex + High first.'

# 6 Endpoint conversions pie
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6), 'Endpoint Conversions', font_size=28, bold=True)
slide.shapes.add_picture(pie_endpoints, Inches(1.5), Inches(1.2), width=Inches(6))
slide.notes_slide.notes_text_frame.text = 'LOV-heavy mapping requires backend lookups; validations need split front/back.'

# 7 Before/After FMB size
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6), 'FMB Size: Before / After', font_size=28, bold=True)
slide.shapes.add_picture(bar_fmb, Inches(1.5), Inches(1.1), width=Inches(6))
add_textbox(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(1), 'Result: ~62% file size reduction and significant noise removal.', font_size=16)
slide.notes_slide.notes_text_frame.text = 'Stress performance and maintainability gains.'

# 8 Business logic & DB interactions (flow text)
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6), 'Business Logic & DB Interactions', font_size=28, bold=True)
flow_text = (
    f"- Heavy business logic consolidated into {METRICS['program_units_procedures']} program-unit procedures.\n"
    f"- Calls to DB packages/procedures/functions were identified (counts tracked).\n"
    f"- 1 report is invoked from forms; requires secure integration.\n"
)
add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(2), flow_text, font_size=16)
slide.notes_slide.notes_text_frame.text = 'Explain separation of concerns and plan to move logic server-side.'

# 9 Security callout
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6), 'Security & Quality Issues (Urgent)', font_size=28, bold=True, color=(200,30,30))
security_text = (
    '- Critical: insecure report call (not authenticated/encrypted)\n'
    '- Hard-coded credentials/values detected\n'
    '- Database connections left open / not securely handled\n'
    '\nSuggested quick wins:\n- Replace hard-coded secrets with vault-backed config\n- Ensure report calls use secured service endpoints\n- Close DB connections and use connection pooling with TLS'
)
add_textbox(slide, Inches(0.6), Inches(1.2), Inches(8.5), Inches(3), security_text, font_size=14, color=(80,80,80))
slide.notes_slide.notes_text_frame.text = 'Prioritize Security fixes in Sprint 0.'

# 10 Deliverables & roadmap
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6), 'Deliverables & Implementation Roadmap', font_size=28, bold=True)
roadmap_text = (
    'Deliverables:\n- Backend endpoint mappings (13 LOV, 6 CRUD, 5 ACTION, 4 VALIDATION)\n- HTML wireframes and Angular FE plan\n\nRoadmap (example):\n- Sprint 0: Security fixes + test harness\n- Sprint 1–2: CRUD & validations\n- Sprint 3: LOV endpoints & FE wiring\n- Sprint 4: Program-unit refactor & DB migration'
)
add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(3.5), roadmap_text, font_size=14)
slide.notes_slide.notes_text_frame.text = 'Map work into sprints and show quick wins.'

# 11 Appendix: Mermaid-like graph (simple text image)
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6), 'Appendix: Logical Graph (Mermaid-like)', font_size=28, bold=True)

# Create a simple PNG with the mermaid-like flow text
mermaid_png = os.path.join(ASSETS_DIR, 'mermaid_sample.png')
fig, ax = plt.subplots(figsize=(8,4))
ax.axis('off')
text = (
    "Form Triggers -> Endpoints -> Program Units -> DB Packages\n"
    "Form A -> calls ProgramUnit.X -> calls DB.Pkg.procY\n"
    "Form B -> LOV endpoints -> Backend -> Table lookups\n"
)
ax.text(0.01, 0.9, text, fontsize=12, va='top')
fig.savefig(mermaid_png, dpi=150, bbox_inches='tight')
plt.close(fig)

slide.shapes.add_picture(mermaid_png, Inches(0.5), Inches(1.0), width=Inches(9))
slide.notes_slide.notes_text_frame.text = 'Provide this diagram as a high-level map; replace with real mermaid export if available.'

# 12 Next steps & closing
slide = prs.slides.add_slide(blank_layout)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6), 'Next Steps & Closing', font_size=28, bold=True)
next_text = (
    '- Approve Sprint 0 to remediate security items\n- Prioritize Complex/High triggers for refactor\n- Start BE endpoints (CRUD & validations) and deliver FE wireframes\n- Measure post-migration KPIs: size, incidents, deploy time\n\nThank you — questions?'
)
add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(3), next_text, font_size=16)
slide.notes_slide.notes_text_frame.text = 'End with a call to action: approve roadmap and secure resources.'

# Save PPTX
prs.save(OUT_PPTX)
print(f'Generated {OUT_PPTX} in the current directory. Run this script locally to produce the PPTX file.')
